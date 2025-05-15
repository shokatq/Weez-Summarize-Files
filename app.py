from flask import Flask, request, jsonify
from flask_cors import CORS
from openai import AzureOpenAI
import os
import tempfile
import logging
import mammoth
import PyPDF2
import pandas as pd
from pptx import Presentation
import uuid

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)

# Temporary storage for extracted text per user
# Key structure: {user_id: {session_id: document_text}}
document_cache = {}

# Conversation history storage
# Key structure: {user_id: {conversation_id: [message1, message2, ...]}}
conversation_cache = {}

def extract_text_from_file(file_path, file_type):
    """Extracts text from different file formats."""
    try:
        logger.info(f"Extracting text from {file_type} file: {file_path}")
        
        if file_type == "pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = "\n".join([page.extract_text() or "" for page in reader.pages])
            return text.strip()

        elif file_type == "docx":
            with open(file_path, "rb") as f:
                raw_text = mammoth.extract_raw_text(f)
            return raw_text.value.strip()

        elif file_type == "pptx":
            presentation = Presentation(file_path)
            text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
            return text.strip()

        elif file_type == "xlsx":
            df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
            text = "\n".join([df[sheet].to_csv(index=False) for sheet in df])
            return text.strip()

        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    except Exception as e:
        logger.error(f"Error extracting text from {file_type}: {str(e)}")
        raise ValueError(f"Failed to extract text from {file_type}: {str(e)}")


@app.route('/summarize', methods=['POST'])
def summarize():
    """Receives a file upload, extracts text, and sends it to OpenAI API."""
    try:
        logger.info("Received summarize request")
        
        # Get user ID from request
        user_id = request.headers.get('X-User-ID')
        if not user_id:
            logger.warning("No user ID provided")
            return jsonify({"error": "User authentication required"}), 401
            
        # Validate request
        if 'file' not in request.files:
            logger.warning("No file uploaded")
            return jsonify({"error": "No file uploaded"}), 400

        uploaded_file = request.files['file']
        file_name = uploaded_file.filename
        
        if not file_name:
            logger.warning("Empty filename")
            return jsonify({"error": "Invalid filename"}), 400
            
        file_type = file_name.split('.')[-1].lower()
        logger.info(f"Processing file: {file_name}, type: {file_type}")

        # Validate file type
        allowed_types = {"pdf", "docx", "pptx", "xlsx"}
        if file_type not in allowed_types:
            logger.warning(f"Unsupported file type: {file_type}")
            return jsonify({"error": f"Unsupported file type: {file_type}"}), 400

        # Save file to a temporary location
        temp_file_path = os.path.join(tempfile.gettempdir(), file_name)
        uploaded_file.save(temp_file_path)
        logger.info(f"File saved to temporary path: {temp_file_path}")

        # Extract text from file
        extracted_text = extract_text_from_file(temp_file_path, file_type)
        logger.info(f"Extracted {len(extracted_text)} characters of text")

        # Ensure valid text extraction
        if not extracted_text.strip():
            logger.warning("No text extracted from file")
            return jsonify({"error": "No text extracted from file"}), 400

        # Generate a new session ID if not provided
        session_id = request.headers.get('X-Session-ID')
        if not session_id:
            session_id = str(uuid.uuid4())
            
        # Store extracted text for follow-up questions, organized by user ID
        if user_id not in document_cache:
            document_cache[user_id] = {}
        document_cache[user_id][session_id] = extracted_text
        logger.info(f"Stored document for user {user_id}, session {session_id}")
        
        # Validate API credentials
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            logger.error("Missing OpenAI API key")
            return jsonify({"error": "Server configuration error: Missing API key"}), 500
            
        # Create client with API key
        client = AzureOpenAI(
            api_key=api_key,
            api_version="2024-12-01-preview",
            azure_endpoint="https://weez-openai-resource.openai.azure.com/"
        )
    
        # Deployment Name (from Azure)
        DEPLOYMENT_NAME = "gpt-4o"  # Change to "gpt-4o" if needed
        logger.info(f"Sending request to OpenAI with model {DEPLOYMENT_NAME}")
        
        response = client.chat.completions.create(
            model=DEPLOYMENT_NAME,
            messages=[
               {"role": "system", "content": """
                You are an expert at explaining complex topics in an engaging and easy-to-understand way. 
                Your task is to summarize the following document in a compelling and professional manner, 
                
                Ensure the summary is **clear, persuasive, and easy to grasp** while maintaining a professional tone.  
                Use **well-placed emojis** to make key points stand out without overloading the text.  
                Here is the text:
                """},
                {"role": "user", "content": extracted_text[:50000]}  # Limit input size
            ],
            temperature=0.3
        )
    
        summary = response.choices[0].message.content
        logger.info("Successfully generated summary")
        
        # Clean up temp file
        try:
            os.remove(temp_file_path)
        except Exception as e:
            logger.warning(f"Failed to remove temp file: {e}")
            
        # Return the session ID along with the summary
        return jsonify({
            "summary": summary,
            "session_id": session_id
        })
    
    except Exception as e:
        logger.error(f"Error in summarize endpoint: {str(e)}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route('/ask', methods=['POST'])
def ask():
    """Handles follow-up questions about the previously uploaded document."""
    try:
        logger.info("Received ask request")
        
        # Get user ID from request
        user_id = request.headers.get('X-User-ID')
        if not user_id:
            logger.warning("No user ID provided")
            return jsonify({"error": "User authentication required"}), 401
            
        session_id = request.headers.get('X-Session-ID')
        if not session_id:
            logger.warning("No session ID provided")
            return jsonify({"error": "Session ID required"}), 400
        
        # Ensure there's a stored document for this user and session
        if user_id not in document_cache or session_id not in document_cache.get(user_id, {}):
            logger.warning(f"No document found for user {user_id}, session {session_id}")
            return jsonify({"error": "No document found. Please upload a file first."}), 400

        # Get the user's query
        data = request.get_json()
        if not data:
            logger.warning("No JSON data in request")
            return jsonify({"error": "Missing request data"}), 400
            
        query = data.get("query")

        if not query or not query.strip():
            logger.warning("Empty query received")
            return jsonify({"error": "Query cannot be empty"}), 400

        # Retrieve the stored document text for this specific user and session
        document_text = document_cache[user_id][session_id]
        logger.info(f"Retrieved document ({len(document_text)} chars) for user {user_id}, query: {query[:50]}...")

        # Validate API credentials
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            logger.error("Missing OpenAI API key")
            return jsonify({"error": "Server configuration error: Missing API key"}), 500

        # Send follow-up question to OpenAI with the document context
        client = AzureOpenAI(
            api_key=api_key,
            api_version="2024-12-01-preview",
            azure_endpoint="https://weez-openai-resource.openai.azure.com/"
        )
    
        # Deployment Name (from Azure)
        DEPLOYMENT_NAME = "gpt-4o"  # Change to "gpt-4o" if needed
        logger.info(f"Sending request to OpenAI with model {DEPLOYMENT_NAME}")
        
        response = client.chat.completions.create(
            model=DEPLOYMENT_NAME,
            messages=[
                {"role": "system", "content": "You are answering questions based on the following document:"},
                {"role": "user", "content": document_text[:15000]},  # Limit input size
                {"role": "user", "content": f"Based on this document, {query}"}
            ],
            temperature=0.3
        )

        answer = response.choices[0].message.content
        logger.info("Successfully generated answer")
        
        return jsonify({"answer": answer})

    except Exception as e:
        logger.error(f"Error in ask endpoint: {str(e)}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route('/chat', methods=['POST'])
def chat():
    """Handles general conversation with Weezy (like ChatGPT)."""
    try:
        logger.info("Received chat request")
        
        # Get user ID from request
        user_id = request.headers.get('X-User-ID')
        if not user_id:
            logger.warning("No user ID provided")
            return jsonify({"error": "User authentication required"}), 401
            
        # Get or create conversation ID
        conversation_id = request.headers.get('X-Conversation-ID')
        if not conversation_id:
            conversation_id = str(uuid.uuid4())
            logger.info(f"Created new conversation {conversation_id} for user {user_id}")
        
        # Get the user's message
        data = request.get_json()
        if not data:
            logger.warning("No JSON data in request")
            return jsonify({"error": "Missing request data"}), 400
            
        message = data.get("message")
        if not message or not message.strip():
            logger.warning("Empty message received")
            return jsonify({"error": "Message cannot be empty"}), 400

        # Initialize conversation history for this user if it doesn't exist
        if user_id not in conversation_cache:
            conversation_cache[user_id] = {}
            
        # Initialize conversation for this conversation ID if it doesn't exist
        if conversation_id not in conversation_cache[user_id]:
            conversation_cache[user_id][conversation_id] = []
        
        # Get conversation history
        conversation_history = conversation_cache[user_id][conversation_id]
        
        # Add user message to history
        conversation_history.append({"role": "user", "content": message})
        
        # Validate API credentials
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            logger.error("Missing OpenAI API key")
            return jsonify({"error": "Server configuration error: Missing API key"}), 500

        # Prepare messages for API
        messages = [
            {"role": "system", "content": """
            You are Weezy, a helpful, friendly, and knowledgeable AI assistant. 
            You provide concise, accurate information in a conversational tone.
            You can discuss a wide range of topics but will politely decline requests that are harmful or inappropriate.
            Use emojis occasionally to keep the conversation engaging.
            """}
        ]
        
        # Add conversation history (limit to last 10 messages to avoid token limits)
        messages.extend(conversation_history[-10:])
        
        # Send to OpenAI
        client = AzureOpenAI(
            api_key=api_key,
            api_version="2024-12-01-preview",
            azure_endpoint="https://weez-openai-resource.openai.azure.com/"
        )
    
        # Deployment Name (from Azure)
        DEPLOYMENT_NAME = "gpt-4o"  # Change to "gpt-4o" if needed
        logger.info(f"Sending request to OpenAI with model {DEPLOYMENT_NAME}")
        
        response = client.chat.completions.create(
            model=DEPLOYMENT_NAME,
            messages=messages,
            temperature=0.7  # Slightly higher temperature for more creative responses
        )

        assistant_response = response.choices[0].message.content
        logger.info(f"Generated response for chat: {assistant_response[:50]}...")
        
        # Add assistant response to conversation history
        conversation_history.append({"role": "assistant", "content": assistant_response})
        
        # Save updated conversation history
        conversation_cache[user_id][conversation_id] = conversation_history
        
        # Return response with conversation ID
        return jsonify({
            "response": assistant_response,
            "conversation_id": conversation_id
        })

    except Exception as e:
        logger.error(f"Error in chat endpoint: {str(e)}", exc_info=True)
        return jsonify({"error": str(e)}), 500


# Add a cleanup endpoint to explicitly remove user data when needed
@app.route('/cleanup', methods=['POST'])
def cleanup():
    """Removes stored documents for a user or specific session."""
    try:
        user_id = request.headers.get('X-User-ID')
        if not user_id:
            return jsonify({"error": "User ID required"}), 400
            
        session_id = request.headers.get('X-Session-ID')
        conversation_id = request.headers.get('X-Conversation-ID')
        data_type = request.args.get('type', 'all')  # 'document', 'conversation', or 'all'
        
        result = {"status": "No data removed"}
        status_code = 200
        
        # Clean up document data
        if data_type in ['document', 'all'] and user_id in document_cache:
            if session_id:
                # Remove specific session
                if session_id in document_cache[user_id]:
                    del document_cache[user_id][session_id]
                    logger.info(f"Removed document session {session_id} for user {user_id}")
                    result["document_status"] = "Session data removed"
                else:
                    result["document_status"] = "Session not found"
            else:
                # Remove all document sessions
                del document_cache[user_id]
                logger.info(f"Removed all document sessions for user {user_id}")
                result["document_status"] = "All document data removed"
        
        # Clean up conversation data
        if data_type in ['conversation', 'all'] and user_id in conversation_cache:
            if conversation_id:
                # Remove specific conversation
                if conversation_id in conversation_cache[user_id]:
                    del conversation_cache[user_id][conversation_id]
                    logger.info(f"Removed conversation {conversation_id} for user {user_id}")
                    result["conversation_status"] = "Conversation removed"
                else:
                    result["conversation_status"] = "Conversation not found"
            else:
                # Remove all conversations
                del conversation_cache[user_id]
                logger.info(f"Removed all conversations for user {user_id}")
                result["conversation_status"] = "All conversations removed"
        
        return jsonify(result), status_code
            
    except Exception as e:
        logger.error(f"Error in cleanup endpoint: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/health', methods=['GET'])
def health_check():
    """Simple endpoint to verify the API is running."""
    return jsonify({"status": "ok"})


if __name__ == '__main__':
    logger.info("Starting Flask server")
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port, debug=False, threaded=True)  # Disable debug in production
